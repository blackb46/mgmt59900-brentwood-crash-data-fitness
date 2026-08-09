"""Generate the script and recording instructions FROM the deployed deck.

The deployed .pptx is the source of truth: Kevin edits it directly, so titles,
speaker notes, and per-slide timings are all read out of it rather than
maintained separately. Nothing here writes to the deck.
"""
import os
import re
import sys

import docx
from docx.shared import Inches, Pt
from pptx import Presentation

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from docx_helpers import body, bullet, h1, h2, header_block, memoline, table

SCR = os.path.dirname(os.path.abspath(__file__))
PRES = (r"C:/Users/kevin/OneDrive - City of Brentwood/Documents/COWORK_MASTER"
        r"/projects/purdue_ai/current_courses/MGMT59900_BigDataCloud"
        r"/Portfolio_Project/Final_Package/Presentation")
DECK = os.path.join(PRES, "Group_15_Kevin_Blackburn_Presentation.pptx")

WPM = 135
TITLE_SLIDE_NAME = "Blind Spots in Public Crash Data (title slide)"

DELIVERY = [
    "Say the reframe out loud inside the first thirty seconds. Do not save it. A "
    "listener who hears early that the data failed will read everything after it "
    "as deliberate rather than as a project that fell short.",
    "This is where your professional credibility shows. You are not a student "
    "describing a dataset, you are the City Engineer describing a budget "
    "constraint. Point at the red line on the map when you say Interstate 65.",
    "Your densest slide and your longest. Do not read the boxes. Say four "
    "sources, name the five layers, then spend what is left on the Redshift "
    "decision, because choosing not to build something is what shows judgment.",
    "The cloud-versus-desktop agreement is the strongest technical moment in the "
    "talk. Land it as one sentence and pause for a beat before moving on.",
    "Four numbers in about seventy seconds: 29.9, 4.4, 87.5, sixteen-fold. Say "
    "each cleanly and do not editorialize between them. Then set up the handoff "
    "to slide 6 deliberately.",
    "The payoff, and it is a story rather than a chart reading. Walk it in order: "
    "the contradiction, why it was the tell, the re-test, the inverted answer. "
    "Then let the two maps sit on screen for a second in silence.",
    "The only purely technical slide, and the easiest place to recover time if "
    "you are running long. The 1.4 second figure is the one that makes the "
    "platform verdict credible rather than defensive.",
    "This decides how the talk is remembered. Slow down. Let the word OUTCOME "
    "land, then RECOMMENDATION, then NEXT STEPS, so the structure is audible and "
    "not only visible. Say thank you and stop talking.",
]


def clock(seconds):
    seconds = int(round(seconds / 5.0) * 5)
    return f"{seconds // 60}:{seconds % 60:02d}"


def split_paragraphs(text, max_sentences=3):
    sents = re.split(r'(?<=[.?!])\s+', ' '.join(text.split()))
    out, cur = [], []
    for s in sents:
        cur.append(s)
        if len(cur) >= max_sentences:
            out.append(' '.join(cur))
            cur = []
    if cur:
        out.append(' '.join(cur))
    return out


def new_doc():
    d = docx.Document()
    st = d.styles["Normal"]
    st.font.name = "Arial"
    st.font.size = Pt(10.5)
    for sec in d.sections:
        sec.left_margin = sec.right_margin = Inches(0.75)
        sec.top_margin = sec.bottom_margin = Inches(0.75)
    return d


# ------------------------------------------------------------- read the deck
prs = Presentation(DECK)
slides = []
cum = 0.0
for i, s in enumerate(prs.slides):
    title = ''
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            title = ph.text_frame.text.strip()
    if not title:
        title = TITLE_SLIDE_NAME
    note = s.notes_slide.notes_text_frame.text if s.has_notes_slide else ''
    words = len(note.split())
    secs = words / WPM * 60
    cum += secs
    slides.append({'n': i + 1, 'title': title, 'notes': note, 'words': words,
                   'time': clock(secs), 'cum': clock(cum)})

total_words = sum(s['words'] for s in slides)
total_secs = total_words / WPM * 60
TOTAL = clock(total_secs)
FAST = clock(total_words / 150 * 60)

# ==================================================================== script
doc = new_doc()
header_block(doc, "MGMT 59900: Big Data Analytics in the Cloud",
             "Final Presentation Script",
             "Kevin Blackburn, P.E., GISP  |  Group 15",
             "Recording due Thursday, August 13, 2026")
body(doc,
     "This script is generated directly from the speaker notes in the deck, so the "
     "two are always identical. Read from Presenter View or from here, whichever "
     "you prefer. If you edit a note in PowerPoint, regenerate this document rather "
     "than editing it by hand.")

h1(doc, "Recording Targets")
table(doc, ["Item", "Target"],
      [["Total runtime", f"{TOTAL} at a normal pace, against a required 8:00 to 10:00"],
       ["Spoken length", f"{total_words:,} words across {len(slides)} slides"],
       ["Format", "Face plus slides, recorded in PowerPoint with Cameo, exported to MP4"],
       ["Deck", "Group_15_Kevin_Blackburn_Presentation.pptx"],
       ["Submission", "Brightspace by 11:59 PM ET, Thursday, August 13, 2026"]],
      widths=[1.6, 5.0])

h1(doc, "Slide Structure")
table(doc, ["#", "Slide", "Words", "Time", "Cumulative"],
      [[str(s['n']), s['title'], str(s['words']), s['time'], s['cum']] for s in slides],
      widths=[0.4, 3.6, 0.7, 0.8, 1.1])
body(doc,
     "Slides 5 and 6 carry the argument and run back to back: 5 shows the dataset "
     "is thin, 6 shows it is also biased and inverts the recommendation. Do not "
     "rush the handoff between them. Slide 8 is the close. If you run long, take "
     "the time out of slides 3 and 7, which are the two a listener can also read "
     "off the deck.")

h1(doc, "Pace")
body(doc,
     f"The spoken script is {total_words:,} words. That is {TOTAL} at a normal "
     f"presenting pace of {WPM} words per minute and {FAST} at a brisk 150. The "
     "risk is at the slow end, where 120 words per minute would run past the "
     "ceiling. If your rehearsal comes in above 9:45, do not cut content. Pick up "
     "the pace slightly and shorten the pauses between slides, which is where most "
     "of the drift accumulates.")

h1(doc, "Per-Slide Script")
for s in slides:
    h2(doc, f"Slide {s['n']}: {s['title']}")
    memoline(doc, "Target time:", f"{s['time']}   ({s['words']} words)")
    memoline(doc, "What to say:", "")
    for para in split_paragraphs(s['notes']):
        body(doc, para)
    memoline(doc, "Delivery note:", DELIVERY[s['n'] - 1])

h1(doc, "Before You Record")
bullet(doc, "Read the whole script aloud once against a timer. Adjust pace, not content.")
bullet(doc, "The deck's Notes pane holds this same text, so Presenter View gives you "
            "everything on screen.")
bullet(doc, "Do not read the on-screen bullets verbatim. The audience reads them "
            "faster than you can say them.")
bullet(doc, "If you fumble a number, stop and re-record that slide only. PowerPoint "
            "records per slide.")
doc.save(os.path.join(SCR, "Blackburn_Presentation_Script.docx"))

# ==================================================== recording instructions
doc = new_doc()
header_block(doc, "MGMT 59900: Big Data Analytics in the Cloud",
             "Final Presentation Recording Instructions",
             "Kevin Blackburn, P.E., GISP  |  Group 15",
             "Recording due Thursday, August 13, 2026")
body(doc,
     f"Deck: Group_15_Kevin_Blackburn_Presentation.pptx, {len(slides)} slides, "
     f"{TOTAL} target against a required 8 to 10 minute window. Speaker notes are "
     "in the Notes pane of every slide and match the script document exactly.")

h1(doc, "Final Slide Order")
table(doc, ["#", "Slide", "Time", "Cumulative"],
      [[str(s['n']), s['title'], s['time'], s['cum']] for s in slides],
      widths=[0.4, 4.2, 1.0, 1.0])

h1(doc, "Step 1. Rehearse Once with a Timer")
bullet(doc, "Open the deck in Presenter View, or keep the script document on a "
            "second monitor.")
bullet(doc, f"Read it aloud once against a stopwatch. Target {TOTAL}.")
bullet(doc, "If you land under 8:00 or over 10:00, adjust pace first. Only cut "
            "content if you are still outside after a second pass.")

h1(doc, "Step 2. Add Cameo So Your Face Appears")
bullet(doc, "On each slide where you want your face visible, go to Insert, then "
            "Cameo. PowerPoint adds a circular camera frame.")
bullet(doc, "Resize and move the circle to a corner that does not cover content. On "
            "slide 3 keep it clear of the architecture diagram, and on slide 6 keep "
            "it clear of the two maps.")
bullet(doc, "Cameo on the title slide and the closing slide is enough if you would "
            "rather not have it on every slide.")

h1(doc, "Step 3. Record Slide by Slide")
bullet(doc, "Go to the Record tab on the ribbon and click Record.")
bullet(doc, "Press the red Record button and wait for the three second countdown.")
bullet(doc, "Talk through the slide using the notes in the pane, then click forward. "
            "Voice, face, and timing are captured together.")
bullet(doc, "To redo one slide: navigate to it, click Clear, then Clear Recordings "
            "on Current Slide, and record it again. You do not restart the deck.")
bullet(doc, f"When all {len(slides)} slides are recorded, save the .pptx. The audio "
            "and video are embedded in the file.")

h1(doc, "Step 4. Export as MP4")
bullet(doc, "File, then Export, then Create a Video.")
bullet(doc, "Set quality to Full HD 1080p and confirm Use Recorded Timings and "
            "Narrations is selected.")
bullet(doc, "Save as Blackburn_MGMT59900_Final_Presentation.mp4 in the Presentation "
            "folder.")
bullet(doc, "Export takes several minutes. Watch the progress bar at the bottom of "
            "the window before closing PowerPoint.")

h1(doc, "Step 5. Check Before You Submit")
bullet(doc, "Play the MP4 start to finish and confirm the runtime is between 8:00 "
            "and 10:00.")
bullet(doc, f"Confirm audio is audible on all {len(slides)} slides and that none is "
            "silent.")
bullet(doc, "Confirm the architecture diagram on slide 3 and the two maps on slide 6 "
            "are readable at full screen.")
bullet(doc, "Upload to Brightspace by 11:59 PM ET on Thursday, August 13, 2026.")

h1(doc, "Note on the Peer Reviews")
body(doc,
     "Two peer reviews are due Saturday, August 15. They are a separate deliverable "
     "and follow the same pattern as the monitor report from ECE 570: a short "
     "summary of what the presenter built, then two or three substantive questions. "
     "Watch the assigned presentations first, then write them.")
doc.save(os.path.join(SCR, "Recording_Instructions.docx"))

print(f"total {total_words} words | {TOTAL} at {WPM} wpm | {FAST} at 150")
for s in slides:
    print(f"  {s['n']}: {s['words']:3d} w  {s['time']}  cum {s['cum']}  {s['title'][:48]}")
